import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE,PP_ALIGN
from scipy.misc import toimage,imsave, imread

## EDA Tools
def add_text(slide,left,top,text,fontsize=14):
    left_dis = left
    top_dis = top
    txtwid = Inches(2)
    txtheig = Inches(0.5)
    txbox = slide.shapes.add_textbox(left_dis, top_dis, txtwid, txtheig)
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(fontsize)
    font.bold = True
    font.color.rgb = RGBColor(0, 0, 255)
def strip_filename(filepath):
    return (filepath.split("/")[-1]).split(".")[0]
def max_side(w, h, max_side=4):
    if w >= h:
        width = Inches(max_side)
        height = width / w * h
    else:
        height = Inches(max_side)
        width = height / h * w
    return width, height
def add_4_pic_slide(prs, pics):
    blank_slide_layout = prs.slide_layouts[5]
    c_left_list = [0.2,5.5,0.2,5.5]
    c_top_list = [0.2,0.2,3.5,3.5]
    slide = prs.slides.add_slide(blank_slide_layout)
    for j in xrange(4):
        img_path = pics[j]
        imgname = strip_filename(img_path)
        im = Image.open(img_path)
        w, h = im.size
        c_left = Inches(c_left_list[j])
        c_top = Inches(c_top_list[j])
        width, height = max_side(w, h, max_side=3.2)
        res_text = imgname + str(w) + " x " + str(h)
        add_text(slide, c_left, c_top - Inches(0.2), res_text, fontsize=6)
        slide.shapes.add_picture(img_path, c_left, c_top, width=width, height=height)

def create_pptx(pptname = "test.pptx", image_list = None):
    prs = Presentation()
    n_pics = len(image_list)
    print(n_pics,n_pics/8)
    for i in xrange(n_pics/8):
        pic_names = [image_list[i * 8 + j * 2] for j in xrange(4)]
        add_4_pic_slide(prs,pic_names)
        pic_names = [image_list[i * 8 + j * 2 + 1] for j in xrange(4)]
        add_4_pic_slide(prs,pic_names)
    prs.save(pptname)

## Image pre-processing
### Gray color conversion
def standard_gamma_corrected(imgarray):
    if np.max(imgarray) > 200: ## If it is a RGB with channel maximum 256
        imgarray = imgarray/255.0
    imgarray[:, :, 0] = imgarray[:, :, 0] ** (1 / 2.2)
    imgarray[:, :, 1] = imgarray[:, :, 1] ** (1 / 2.2)
    imgarray[:, :, 2] = imgarray[:, :, 2] ** (1 / 2.2)
    return imgarray
def gleam_rgb2gray(imgpath, imgname, output_folder = "./"):
    imgarray = imread(imgpath)
    if np.ndim(imgarray) == 3:
        gc = standard_gamma_corrected(imgarray)
        gray = 1/3. * (gc[:, :, 0] + gc[:, :, 1] + gc[:, :, 2])
        imsave(output_folder + imgname[0:-4] + ".png", gray)
    else:
        imsave(output_folder + imgname[0:-4] + ".png", imgarray)
    return gray
def luminance_rgb2gray(imgpath, imgname, output_folder = "./"):
    imgarray = imread(imgpath)
    if np.ndim(imgarray) == 3:
        gray = imgarray[:, :, 0] * 0.3 + imgarray[:, :, 1] * 0.59 + imgarray[:, :, 2] * 0.11
        imsave(output_folder + imgname[0:-4] + ".png", gray)
    else:
        imsave(output_folder + imgname[0:-4] + ".png", imgarray)
    return gray

